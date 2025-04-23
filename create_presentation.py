from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_title_slide(prs):
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = "CareerPortal"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 70, 193)
    
    subtitle.text = "A Modern Job Portal Platform\nBuilt with React, Node.js, and MongoDB"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)

def create_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Project Overview"
    
    overview_text = """
    • Purpose: Connect job seekers with tech companies
    
    Target Users:
    • Job Seekers in Tech Industry
    • Recruiters and Companies
    • HR Administrators
    """
    content.text = overview_text

def create_features_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Key Features"
    
    features_text = """
    For Job Seekers:
    • Easy job search and application
    • Profile management
    • Application tracking
    • Resume submission

    For Employers/Admin:
    • Job posting management
    • Application review system
    • Interview scheduling
    • Candidate status tracking
    """
    content.text = features_text

def create_tech_stack_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technology Stack"
    
    tech_text = """
    Frontend:
    • React.js
    • TailwindCSS
    • React Router
    • Axios

    Backend:
    • Node.js
    • Express.js
    • MongoDB
    • JWT Authentication
    """
    content.text = tech_text

def create_ui_slides(prs):
    # Landing Page
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "User Interface - Landing Page"
    content.text = """
    • Modern, responsive design
    • Engaging hero section
    • Feature highlights
    • Intuitive navigation
    """
    
    # Add more UI slides here...

def create_implementation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Key Implementation Features"
    
    impl_text = """
    Authentication & Security:
    • JWT-based authentication
    • Secure password handling
    • Protected routes

    Database Design:
    • MongoDB collections
    • Efficient data relationships
    • Scalable structure
    """
    content.text = impl_text

def create_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Project Architecture"
    
    arch_text = """
    careers-portal/
    ├── frontend/
    │   ├── src/
    │   │   ├── components/
    │   │   ├── pages/
    │   │   └── services/
    └── backend/
        ├── controllers/
        ├── models/
        ├── routes/
        └── middleware/
    """
    content.text = arch_text

def create_future_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Future Enhancements"
    
    future_text = """
    Planned Features:
    • AI-powered job matching
    • Advanced search filters
    • Email notifications
    • Mobile app version

    Scalability Plans:
    • Cloud deployment
    • Performance optimization
    • Enhanced security
    """
    content.text = future_text

def create_thank_you_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Thank You!"
    subtitle.text = "Questions & Answers"

def main():
    prs = Presentation()
    
    # Create all slides
    create_title_slide(prs)
    create_overview_slide(prs)
    create_features_slide(prs)
    create_tech_stack_slide(prs)
    create_ui_slides(prs)
    create_implementation_slide(prs)
    create_architecture_slide(prs)
    create_future_slide(prs)
    create_thank_you_slide(prs)
    
    # Save the presentation
    prs.save('CareerPortal_Presentation.pptx')
    print("Presentation created successfully!")

if __name__ == "__main__":
    main() 